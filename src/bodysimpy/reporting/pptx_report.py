from pathlib import Path

from pptx import Presentation as create_presentation
from pptx.presentation import Presentation

from bodysimpy.reporting.summary_data import ReportingSummary


def _value(
    value: float | None,
    suffix: str = "",
) -> str:
    if value is None:
        return "Not available"

    return f"{value:.3f}{suffix}"


def _add_title_slide(
    presentation: Presentation,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])

    slide.shapes.title.text = "BodySimPy Management Summary"

    slide.placeholders[1].text = "Python-driven CAE workflow automation"


def _add_bullets(
    presentation: Presentation,
    title: str,
    bullets: list[str],
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])

    slide.shapes.title.text = title

    body = slide.placeholders[1].text_frame
    body.clear()

    for index, bullet in enumerate(bullets):
        if index == 0:
            paragraph = body.paragraphs[0]
        else:
            paragraph = body.add_paragraph()

        paragraph.text = bullet
        paragraph.level = 0


def generate_management_summary_pptx(
    summary: ReportingSummary,
    output_path: str | Path,
) -> None:
    path = Path(output_path)

    path.parent.mkdir(
        parents=True,
        exist_ok=True,
    )

    presentation = create_presentation()

    _add_title_slide(presentation)

    _add_bullets(
        presentation,
        "Objective",
        [
            ("Automate structural CAE workflows using Python."),
            ("Validate FEA against analytical references."),
            ("Generate reusable engineering evidence for portfolio discussion."),
        ],
    )

    _add_bullets(
        presentation,
        "Method",
        [
            "Validated YAML configuration.",
            "Analytical cantilever reference solution.",
            "CalculiX static and modal simulations.",
            ("Stochastic, fatigue and ML surrogate studies."),
        ],
    )

    _add_bullets(
        presentation,
        "Key Results",
        [
            (f"Mesh displacement error: {_value(summary.baseline_deflection_error_percent, ' %')}"),
            (f"Mesh stress error: {_value(summary.baseline_stress_error_percent, ' %')}"),
            (f"Mode-1 frequency: {_value(summary.mode_1_frequency_hz, ' Hz')}"),
            (f"Stress exceedance over 350 MPa: {_value(summary.stress_exceedance_percent, ' %')}"),
            (f"Best ML frequency MAPE: {_value(summary.best_ml_frequency_mape_percent, ' %')}"),
        ],
    )

    _add_bullets(
        presentation,
        "Sensitivity",
        [
            ("Wall-thickness influence on stress, deflection, mass and modal response."),
            ("Material stiffness and density influence on structural dynamics."),
            ("Stochastic uncertainty propagation for manufacturing, material and load inputs."),
        ],
    )

    _add_bullets(
        presentation,
        "Risk & Limitations",
        [
            ("Simplified structural surrogate, not production body-in-white geometry."),
            ("Idealized boundary conditions and generic material assumptions."),
            ("ML surrogate is valid only inside the sampled engineering design space."),
        ],
    )

    _add_bullets(
        presentation,
        "Engineering Recommendation",
        [
            ("Use BodySimPy as a reproducible CAE automation and validation workflow."),
            ("Use automated parameter studies to identify influential design variables."),
            (
                "Extend shell modelling, joint modelling "
                "and durability fidelity before "
                "production-style conclusions."
            ),
        ],
    )

    presentation.save(str(path))
